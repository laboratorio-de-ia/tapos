#!/usr/bin/env python3

"""
===============================================================================
Projeto:
Descrição:

Desenvolvido por Eduardo Ferreira
Cargo: CTO da TAPOS

© 2026 TAPOS. Todos os direitos reservados.
===============================================================================
"""
"""
Conversor Universal de Documentos para Markdown
Suporta: PDF, DOCX, XLSX, CSV, PPTX, TXT, Imagens
Otimizado para economizar tokens ao enviar para Claude
"""

import os
import sys
import pandas as pd
from pathlib import Path
from typing import Optional, List
import json
from datetime import datetime

# ==================== CONFIGURAÇÕES ====================

class ConfiguradorDependencias:
    """Gerencia instalação de dependências"""
    
    @staticmethod
    def instalar_se_necessario():
        """Instala dependências necessárias"""
        dependencias = {
            'pdf2image': 'pdf2image',
            'PIL': 'Pillow',
            'pytesseract': 'pytesseract',
            'python_docx': 'python-docx',
            'pptx': 'python-pptx',
            'pandas': 'pandas',
            'requests': 'requests',
        }
        
        print("🔍 Verificando dependências...")
        faltantes = []
        
        for modulo, pacote in dependencias.items():
            try:
                __import__(modulo)
                print(f"  ✅ {pacote}")
            except ImportError:
                print(f"  ⚠️  {pacote} - será instalado")
                faltantes.append(pacote)
        
        if faltantes:
            print(f"\n📦 Instalando: {', '.join(faltantes)}")
            os.system(f"pip install {' '.join(faltantes)} --break-system-packages")
            print("✅ Dependências instaladas!\n")


# ==================== CONVERSORES ====================

class ConversorPDF:
    """Converte PDF para Markdown"""
    
    @staticmethod
    def converter(caminho_pdf: str, extrair_imagens: bool = False) -> str:
        """
        Converte PDF para Markdown
        
        Args:
            caminho_pdf: Caminho do arquivo PDF
            extrair_imagens: Se True, extrai imagens do PDF
        
        Returns:
            String com conteúdo em Markdown
        """
        try:
            try:
                # Tenta usar pdfplumber (melhor para tabelas)
                import pdfplumber
                print(f"📄 Convertendo PDF com pdfplumber: {caminho_pdf}")
                
                markdown = []
                with pdfplumber.open(caminho_pdf) as pdf:
                    for i, page in enumerate(pdf.pages, 1):
                        markdown.append(f"## Página {i}\n")
                        
                        # Extrai texto
                        texto = page.extract_text()
                        if texto:
                            markdown.append(texto)
                        
                        # Extrai tabelas
                        tabelas = page.extract_tables()
                        if tabelas:
                            for tabela in tabelas:
                                df = pd.DataFrame(tabela[1:], columns=tabela[0])
                                markdown.append(df.to_markdown(index=False))
                        
                        markdown.append("\n---\n")
                
                return "\n".join(markdown)
            
            except ImportError:
                # Fallback: usa pytesseract para OCR
                print(f"📄 Convertendo PDF com OCR: {caminho_pdf}")
                from pdf2image import convert_from_path
                from pytesseract import image_to_string
                
                imagens = convert_from_path(caminho_pdf)
                markdown = []
                
                for i, imagem in enumerate(imagens, 1):
                    markdown.append(f"## Página {i}\n")
                    texto = image_to_string(imagem, lang='por')
                    markdown.append(texto)
                    markdown.append("\n---\n")
                
                return "\n".join(markdown)
        
        except Exception as e:
            print(f"❌ Erro ao converter PDF: {e}")
            return f"# Erro ao converter PDF\n\n{str(e)}"


class ConversorDocx:
    """Converte Word (.docx) para Markdown"""
    
    @staticmethod
    def converter(caminho_docx: str) -> str:
        """Converte DOCX para Markdown preservando estrutura"""
        try:
            from docx import Document

            print(f"📝 Convertendo DOCX: {caminho_docx}")
            doc = Document(caminho_docx)
            markdown = []
            
            for para in doc.paragraphs:
                if para.text.strip():
                    # Detecta nível de título
                    if 'Heading 1' in para.style.name:
                        markdown.append(f"# {para.text}\n")
                    elif 'Heading 2' in para.style.name:
                        markdown.append(f"## {para.text}\n")
                    elif 'Heading 3' in para.style.name:
                        markdown.append(f"### {para.text}\n")
                    elif para.style.name.startswith('List'):
                        markdown.append(f"- {para.text}\n")
                    else:
                        markdown.append(f"{para.text}\n")
            
            # Extrai tabelas
            for tabela in doc.tables:
                dados = []
                for linha in tabela.rows:
                    dados.append([celula.text for celula in linha.cells])
                
                if dados:
                    df = pd.DataFrame(dados[1:], columns=dados[0])
                    markdown.append(df.to_markdown(index=False))
                    markdown.append("\n")
            
            return "\n".join(markdown)
        
        except Exception as e:
            print(f"❌ Erro ao converter DOCX: {e}")
            return f"# Erro ao converter DOCX\n\n{str(e)}"


class ConversorExcel:
    """Converte Excel/CSV para Markdown"""
    
    @staticmethod
    def converter(caminho_excel: str, sheet: Optional[str] = None) -> str:
        """Converte XLSX/CSV para Markdown"""
        try:
            print(f"📊 Convertendo Excel: {caminho_excel}")
            
            if caminho_excel.endswith('.csv'):
                df = pd.read_csv(caminho_excel)
                markdown = df.to_markdown(index=False)
            else:
                # Lê todas as abas
                excel_file = pd.ExcelFile(caminho_excel)
                markdown = []
                
                for nome_aba in excel_file.sheet_names:
                    markdown.append(f"## Aba: {nome_aba}\n")
                    df = pd.read_excel(caminho_excel, sheet_name=nome_aba)
                    markdown.append(df.to_markdown(index=False))
                    markdown.append("\n")
                
                markdown = "\n".join(markdown)
            
            return markdown
        
        except Exception as e:
            print(f"❌ Erro ao converter Excel: {e}")
            return f"# Erro ao converter Excel\n\n{str(e)}"


class ConversorPowerPoint:
    """Converte PowerPoint para Markdown"""
    
    @staticmethod
    def converter(caminho_pptx: str) -> str:
        """Converte PPTX para Markdown"""
        try:
            from pptx import Presentation
            
            print(f"📽️  Convertendo PowerPoint: {caminho_pptx}")
            prs = Presentation(caminho_pptx)
            markdown = []
            
            for i, slide in enumerate(prs.slides, 1):
                markdown.append(f"## Slide {i}\n")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        markdown.append(f"{shape.text}\n")
                
                markdown.append("\n---\n")
            
            return "\n".join(markdown)
        
        except Exception as e:
            print(f"❌ Erro ao converter PowerPoint: {e}")
            return f"# Erro ao converter PowerPoint\n\n{str(e)}"


class ConversorImagem:
    """Converte imagens para Markdown com OCR"""
    
    @staticmethod
    def converter(caminho_imagem: str) -> str:
        """Extrai texto de imagem usando OCR"""
        try:
            from PIL import Image
            from pytesseract import image_to_string
            
            print(f"🖼️  Extraindo texto de imagem: {caminho_imagem}")
            img = Image.open(caminho_imagem)
            texto = image_to_string(img, lang='por')
            
            return f"# Texto extraído de {Path(caminho_imagem).name}\n\n{texto}"
        
        except Exception as e:
            print(f"❌ Erro ao converter imagem: {e}")
            return f"# Erro ao converter imagem\n\n{str(e)}"


class ConversorTexto:
    """Converte texto plano (.txt) para Markdown"""

    @staticmethod
    def converter(caminho_txt: str) -> str:
        """Lê o texto e devolve como Markdown (conversão direta)"""
        try:
            print(f"📄 Convertendo TXT: {caminho_txt}")

            with open(caminho_txt, 'r', encoding='utf-8', errors='replace') as f:
                return f.read()

        except Exception as e:
            print(f"❌ Erro ao converter TXT: {e}")
            return f"# Erro ao converter TXT\n\n{str(e)}"


class ConversorMarkdown:
    """Repassa arquivos Markdown (.md) sem alterações"""

    @staticmethod
    def converter(caminho_md: str) -> str:
        """Lê o Markdown e devolve como está (conversão direta)"""
        try:
            print(f"📄 Convertendo MD: {caminho_md}")

            with open(caminho_md, 'r', encoding='utf-8', errors='replace') as f:
                return f.read()

        except Exception as e:
            print(f"❌ Erro ao converter MD: {e}")
            return f"# Erro ao converter Markdown\n\n{str(e)}"


# ==================== CONVERSOR PRINCIPAL ====================

class ConvertorUniversal:
    """Conversor automático para qualquer formato"""
    
    CONVERSORES = {
        '.pdf': ConversorPDF.converter,
        '.docx': ConversorDocx.converter,
        '.doc': ConversorDocx.converter,
        '.xlsx': ConversorExcel.converter,
        '.xls': ConversorExcel.converter,
        '.csv': ConversorExcel.converter,
        '.pptx': ConversorPowerPoint.converter,
        '.ppt': ConversorPowerPoint.converter,
        '.png': ConversorImagem.converter,
        '.jpg': ConversorImagem.converter,
        '.jpeg': ConversorImagem.converter,
        '.txt': ConversorTexto.converter,
        '.md': ConversorMarkdown.converter,
    }
    
    def __init__(self, diretorio_saida: str = "./markdown_saida"):
        self.diretorio_saida = Path(diretorio_saida)
        self.diretorio_saida.mkdir(exist_ok=True)
        self.relatorio = {
            'data': datetime.now().isoformat(),
            'arquivos_processados': [],
            'erros': []
        }
    
    def converter_arquivo(self, caminho_arquivo: str) -> Optional[str]:
        """Converte um único arquivo"""
        caminho = Path(caminho_arquivo)
        
        if not caminho.exists():
            print(f"❌ Arquivo não encontrado: {caminho_arquivo}")
            self.relatorio['erros'].append(f"Arquivo não encontrado: {caminho_arquivo}")
            return None
        
        extensao = caminho.suffix.lower()
        
        if extensao not in self.CONVERSORES:
            print(f"❌ Formato não suportado: {extensao}")
            self.relatorio['erros'].append(f"Formato não suportado: {extensao}")
            return None
        
        try:
            print(f"\n{'='*60}")
            print(f"🔄 Convertendo: {caminho.name}")
            print(f"{'='*60}")
            
            conversor = self.CONVERSORES[extensao]
            conteudo_markdown = conversor(caminho_arquivo)
            
            # Salva arquivo de saída
            nome_saida = self.diretorio_saida / f"{caminho.stem}.md"
            with open(nome_saida, 'w', encoding='utf-8') as f:
                f.write(conteudo_markdown)
            
            tamanho_original = caminho.stat().st_size / 1024  # KB
            tamanho_convertido = nome_saida.stat().st_size / 1024  # KB
            economia = ((tamanho_original - tamanho_convertido) / tamanho_original) * 100
            
            info = {
                'arquivo_original': caminho.name,
                'arquivo_convertido': nome_saida.name,
                'tamanho_original_kb': round(tamanho_original, 2),
                'tamanho_convertido_kb': round(tamanho_convertido, 2),
                'economia_percentual': round(economia, 2)
            }
            
            self.relatorio['arquivos_processados'].append(info)
            
            print(f"✅ Convertido com sucesso!")
            print(f"   Tamanho: {tamanho_original:.2f} KB → {tamanho_convertido:.2f} KB")
            print(f"   Economia: {economia:.1f}% 💰")
            print(f"   Salvo em: {nome_saida}")
            
            return str(nome_saida)
        
        except Exception as e:
            erro_msg = f"Erro ao converter {caminho.name}: {str(e)}"
            print(f"❌ {erro_msg}")
            self.relatorio['erros'].append(erro_msg)
            return None
    
    def converter_multiplos(self, arquivos: List[str]) -> List[str]:
        """Converte vários arquivos"""
        resultados = []
        for arquivo in arquivos:
            resultado = self.converter_arquivo(arquivo)
            if resultado:
                resultados.append(resultado)
        
        return resultados
    
    def converter_pasta(self, caminho_pasta: str) -> List[str]:
        """Converte todos os arquivos suportados de uma pasta"""
        pasta = Path(caminho_pasta)
        
        if not pasta.is_dir():
            print(f"❌ Pasta não encontrada: {caminho_pasta}")
            return []
        
        extensoes_suportadas = list(self.CONVERSORES.keys())
        arquivos = [f for f in pasta.iterdir() 
                   if f.suffix.lower() in extensoes_suportadas]
        
        print(f"\n📁 Encontrados {len(arquivos)} arquivos na pasta")
        
        return self.converter_multiplos([str(f) for f in arquivos])
    
    def gerar_relatorio(self) -> str:
        """Gera relatório de conversão"""
        relatorio_caminho = self.diretorio_saida / "relatorio.json"
        
        with open(relatorio_caminho, 'w', encoding='utf-8') as f:
            json.dump(self.relatorio, f, indent=2, ensure_ascii=False)
        
        # Imprimir resumo
        print(f"\n{'='*60}")
        print("📊 RELATÓRIO DE CONVERSÃO")
        print(f"{'='*60}")
        print(f"Data: {self.relatorio['data']}")
        print(f"Arquivos processados: {len(self.relatorio['arquivos_processados'])}")
        print(f"Erros: {len(self.relatorio['erros'])}")
        
        if self.relatorio['arquivos_processados']:
            economia_total = sum(a['economia_percentual'] for a in self.relatorio['arquivos_processados'])
            economia_media = economia_total / len(self.relatorio['arquivos_processados'])
            print(f"Economia média: {economia_media:.1f}% 💰")
        
        print(f"Relatório salvo em: {relatorio_caminho}\n")
        
        return str(relatorio_caminho)


# ==================== EXEMPLOS DE USO ====================

def exemplo_uso():
    """Exemplos práticos de uso"""
    
    print("""
    ╔════════════════════════════════════════════════════════════╗
    ║     CONVERSOR UNIVERSAL PYTHON → MARKDOWN                  ║
    ║     Otimizado para economizar tokens com Claude             ║
    ╚════════════════════════════════════════════════════════════╝
    
    EXEMPLOS DE USO:
    
    1. Converter um arquivo:
       ---
       conversor = ConvertorUniversal()
       resultado = conversor.converter_arquivo("documento.pdf")
       
    2. Converter múltiplos arquivos:
       ---
       conversor = ConvertorUniversal()
       resultados = conversor.converter_multiplos([
           "arquivo1.docx",
           "arquivo2.xlsx",
           "arquivo3.pdf"
       ])
    
    3. Converter toda uma pasta:
       ---
       conversor = ConvertorUniversal()
       resultados = conversor.converter_pasta("./meus_documentos")
       conversor.gerar_relatorio()
    
    4. Usar conversores individuais:
       ---
       texto = ConversorPDF.converter("documento.pdf")
       texto = ConversorExcel.converter("dados.xlsx")
       texto = ConversorDocx.converter("relatorio.docx")
    
    FORMATOS SUPORTADOS:
    • PDF (com OCR se necessário)
    • Word (.docx, .doc)
    • Excel (.xlsx, .xls, .csv)
    • PowerPoint (.pptx)
    • Imagens (.png, .jpg, .jpeg)
    """)


if __name__ == "__main__":
    # Instala dependências
    ConfiguradorDependencias.instalar_se_necessario()
    
    # Mostra exemplos
    exemplo_uso()
    
    # Se passar arquivo como argumento
    if len(sys.argv) > 1:
        conversor = ConvertorUniversal()
        
        if Path(sys.argv[1]).is_dir():
            conversor.converter_pasta(sys.argv[1])
        else:
            conversor.converter_arquivo(sys.argv[1])
        
        conversor.gerar_relatorio()
